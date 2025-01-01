"""
File:       PowerPointGenerator.py
Purpose:    This file PowerPointGenerator.py contains the PowerPointGenerator class, which is responsible for the
            creation of the PowerPoint (slides, charts, graphs).
Author:     Joey Borrelli, Software & Training Intern For Catholic Charities of East Tennessee
Anno:       Anno Domini 2024
"""

import pptx as pp
from pptx.chart.data import CategoryChartData
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Pt
import CLI as UI

class PowerPointGenerator:
    # Slide layout constants
    TITLE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3 # side by side bullet textboxes
    COMPARISON = 4 # side by side bullet textboxes with titles
    TITLE_ONLY = 5
    BLANK = 6
    CONTENT_CAPTION = 7
    PICTURE_CAPTION = 8

    def __init__(self):
        # create presentation object
        self.prs = pp.Presentation()

        # create title slide
        slide_layout = self.prs.slide_layouts[self.TITLE]
        slide = self.prs.slides.add_slide(slide_layout)

        # put stuff on the slide
        placeholders = slide.placeholders

        title = placeholders[0]
        subtitle = placeholders[1]

        title.text = UI.get_PowerPoint_Name()
        subtitle.text = "Catholic Charities of East Tennessee"

    def create_Table_Slide(self, title, matrix, columns, rows):
        # check if the numbers given are valid
        if rows > 0 and columns > 0:
            # create a slide
            slide_layout = self.prs.slide_layouts[self.TITLE_ONLY]
            slide = self.prs.slides.add_slide(slide_layout)
            # give the slide a title
            shapes = slide.shapes
            shapes.title.text = title

            # create variables that will be used to position the table
            left = Inches(0.0)
            top = Inches(2.0)
            width = Inches(10.0)
            height = Inches(0.8)

            # create the table
            table = shapes.add_table(rows, columns, left, top, width, height).table

            # fill in cells with data
            for row in range(len(matrix)): # loop through rows
                for col in range(len(matrix[row])): # loop through columns
                    if matrix[row][col] == 'Count':
                        matrix[row][col] = ''
                    table.cell(row, col).text = matrix[row][col]

            # Size the table according to the number of rows, columns, and the longest word
            # Default font size
            font_size = Pt(18)
            # find the longest word
            longest_word_length = 0
            for row in matrix:
                for cell in row:
                    words = cell.split()
                    for word in words:
                        longest_word_length = max(longest_word_length, len(word))
            # set font based on the longest word
            if longest_word_length > 20:
                font_size = Pt(8)
            elif longest_word_length > 15:
                font_size = Pt(10)
            elif longest_word_length > 10:
                font_size = Pt(12)
            # Adjust font size based on number of rows
            if rows > 19:
                font_size = Pt(max(font_size.pt - 2.25, 6))
            elif rows > 15:
                font_size = Pt(max(font_size.pt - 1, 6))
            # Adjust font size based on number of columns
            if columns > 10:
                font_size = Pt(max(font_size.pt - 2, 5))  # Reduce font for many columns, not less than 6pt
            elif columns > 8:
                font_size = Pt(max(font_size.pt - 2, 6))  # Reduce font for many columns, not less than 6pt
            # Apply the calculated font size to all cells in the table
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = font_size

            # merge the first row's cells
            table.cell(0, 0).merge(table.cell(0, columns - 1))
            # center text in merged cell
            for paragraph in table.cell(0, 0).text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER  # Horizontal alignment
            table.cell(0, 0).vertical_alignment = "middle"  # Vertical alignment

            # merge any row's cells where the first element isn't '', but every element after is
            for row in range(len(matrix)):
                if matrix[row][0] != '' and all(cell == '' for cell in matrix[row][1:]):
                    table.cell(row, 0).merge(table.cell(row, columns - 1))
        else:
            print("\nError creating slide " + title + ", rows or columns are < 0")

    def create_PieChart_Slide(self, title, matrix):
        # Create a slide
        slide_layout = self.prs.slide_layouts[self.TITLE_ONLY]
        slide = self.prs.slides.add_slide(slide_layout)
        # give the slide a title
        shapes = slide.shapes
        shapes.title.text = title

        # create a chart
        chart_data = CategoryChartData()
        chart_data.categories = []
        chart_data.add_series('first_series_of_data', (0.25, 0.25, 0.25, 0.25))

        # add chart to slide
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y , cx, cy, chart_data
        ).chart

        # set chart legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_ing_layout = False

        # set chart data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    def create_BarChart_slide(self, title, matrix):
        # create a slide
        slide_layout = self.prs.slide_layouts[self.TITLE_ONLY]
        slide = self.prs.slides.add_slide(slide_layout)
        # give the slide a title
        shapes = slide.shapes
        shapes.title.text = title

        # create a chart
        chart_data = CategoryChartData()
        # fill in x axis (categories)
        chart_data.categories = []
        # fill in with data (you can add multiple series)
        chart_data.add_series('first_series_of_data', (5, 10, 15))

        # add chart to slide
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y , cx, cy, chart_data
        )

    def save_Presentation(self):
        # save the presentation
        self.prs.save("pptx_exports/" + UI.get_PowerPoint_SaveName() + ".pptx")